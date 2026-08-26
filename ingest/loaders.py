"""Local-file ingestion — supplements the URL-based fetchers in sources.py.
Supports the file types most literature actually shows up as: PDF, plain
text/Markdown, HTML, Word, Excel, and PowerPoint documents."""

import glob
from pathlib import Path

from bs4 import BeautifulSoup
from pypdf import PdfReader


def _load_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    text = "\n".join(page.extract_text() or "" for page in reader.pages)
    return text.replace(
        "\x00", ""
    )  # malformed PDFs can yield NUL bytes; Postgres text columns reject them


def _load_text(path: Path) -> str:
    return path.read_text(errors="ignore")


def _load_html(path: Path) -> str:
    soup = BeautifulSoup(path.read_text(errors="ignore"), "html.parser")
    for tag in soup(["script", "style"]):
        tag.decompose()
    return soup.get_text(separator="\n")


def _load_docx(path: Path) -> str:
    from docx import Document

    doc = Document(str(path))
    return "\n".join(p.text for p in doc.paragraphs)


def _load_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    lines = []
    for sheet in wb.worksheets:
        lines.append(f"# {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                lines.append("\t".join(cells))
    return "\n".join(lines)


def _load_pptx(path: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(path))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        lines.append(f"# Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(shape.text_frame.text)
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            lines.append(f"(notes) {slide.notes_slide.notes_text_frame.text}")
    return "\n".join(lines)


LOADERS = {
    ".pdf": _load_pdf,
    ".txt": _load_text,
    ".md": _load_text,
    ".markdown": _load_text,
    ".html": _load_html,
    ".htm": _load_html,
    ".docx": _load_docx,
    ".xlsx": _load_xlsx,
    ".xlsm": _load_xlsx,
    ".pptx": _load_pptx,
}


def load_file(path: str) -> dict | None:
    p = Path(path)
    loader = LOADERS.get(p.suffix.lower())
    if loader is None:
        raise ValueError(f"unsupported file type: {p.suffix} (supported: {', '.join(LOADERS)})")
    text = loader(p)
    if not text or not text.strip():
        return None
    return {
        "source": "file",
        "url": f"file://{p.resolve()}",
        "title": p.stem,
        "text": text,
        "license": "local file — verify you have the right to ingest/redistribute this",
    }


def fetch_local_files(pattern: str) -> list[dict]:
    """`pattern` is a glob, e.g. "./docs/**/*.pdf" (use --path with cli.py)."""
    out = []
    for path in sorted(glob.glob(pattern, recursive=True)):
        p = Path(path)
        if not p.is_file():
            continue
        try:
            doc = load_file(path)
        except Exception as e:
            print(f"  skipping {path}: {type(e).__name__}: {e}")
            continue
        if doc:
            out.append(doc)
    return out
